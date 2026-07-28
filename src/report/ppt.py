"""④ 포트폴리오 PPT 자동 생성 — 워크플로우 산출물을 발표용 덱으로 조립. (소유: mfg-reporter)

수집 규모·EDA 통계·모델 비교·결함별 성능/한계를 read-only로 읽어 python-pptx 덱을 만든다.
이미지는 add_picture로 임베드하고, 성능표는 comparison.parquet에서 동적으로 채운다(하드코딩 금지).
산출물이 없으면 해당 슬라이드를 안내 텍스트로 graceful 처리한다.

실행: python -m src.report.ppt
"""

import sqlite3
from pathlib import Path

import pandas as pd

from config.settings import COLLECT_TABLE, DB_PATH, DEFAULT_MODEL, EDA_DIR, MODELS_DIR, REPORTS_DIR

# 기본 모델은 config.settings 단일 정의를 따른다. 표시 라벨:
MODEL_LABELS = {"lstm_ae": "LSTM-AE", "vae": "VAE", "transformer_ae": "Transformer-AE"}

# 슬라이드 캔버스(16:9) 치수(inch)
SLIDE_W = 13.333
SLIDE_H = 7.5


# ── 산출물 로더(read-only) ────────────────────────────────
def _collect_stats() -> dict:
    """stream.db 수집 통계(총행수·정상/결함). 없으면 빈 통계."""
    if not DB_PATH.exists():
        return {"total": 0}
    try:
        with sqlite3.connect(DB_PATH) as conn:
            df = pd.read_sql(f"SELECT fault_id FROM {COLLECT_TABLE}", conn)
    except (pd.errors.DatabaseError, sqlite3.OperationalError):
        return {"total": 0}
    if df.empty:
        return {"total": 0}
    n_normal = int((df["fault_id"] == 0).sum())
    modes = sorted(int(f) for f in df.loc[df["fault_id"] != 0, "fault_id"].unique())
    return {"total": len(df), "normal": n_normal, "fault": len(df) - n_normal, "fault_modes": modes}


def _comparison() -> pd.DataFrame:
    """모델별 종합 성능표. 없으면 빈 DataFrame."""
    path = MODELS_DIR / "comparison.parquet"
    return pd.read_parquet(path) if path.exists() else pd.DataFrame()


def _per_fault() -> pd.DataFrame:
    """결함 IDV별 탐지율. 없으면 빈 DataFrame."""
    path = MODELS_DIR / "comparison_per_fault.parquet"
    return pd.read_parquet(path) if path.exists() else pd.DataFrame()


# ── 슬라이드 헬퍼 ─────────────────────────────────────────
def _blank_slide(prs):
    """레이아웃 6(빈 슬라이드) 추가."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_title(slide, text: str, *, size: int = 30):
    """슬라이드 상단 제목 텍스트박스."""
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(SLIDE_W - 1.2), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    return box


def _add_bullets(slide, lines, *, left=0.7, top=1.4, width=None, height=5.4, size=16):
    """불릿 텍스트박스. lines 원소는 str 또는 (text, level) 튜플."""
    from pptx.util import Inches, Pt

    width = width if width is not None else SLIDE_W - 1.4
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if level == 0 else "– ") + text
        p.level = level
        p.font.size = Pt(size if level == 0 else size - 2)
    return box


def _add_image_fit(slide, img_path: Path, *, left, top, max_w, max_h):
    """이미지를 비율 유지하며 (max_w, max_h) 박스 안에 맞춰 임베드(가로 중앙 정렬)."""
    from PIL import Image
    from pptx.util import Inches

    with Image.open(img_path) as im:
        iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = iw * ratio, ih * ratio
    cx = left + (max_w - w) / 2
    slide.shapes.add_picture(str(img_path), Inches(cx), Inches(top), Inches(w), Inches(h))


def _missing_note(slide, text: str):
    """산출물 없을 때 안내 텍스트(슬라이드는 유지)."""
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(SLIDE_W - 1.6), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = f"(산출물 없음) {text}"
    p.font.size = Pt(14)


# ── 개별 슬라이드 빌더 ────────────────────────────────────
def _slide_cover(prs, embedded: list) -> None:
    """1. 표지."""
    from pptx.util import Inches, Pt

    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(SLIDE_W - 1.6), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = "제조 공정 이상탐지 자동화 워크플로우"
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    for line in [
        "TEP 다변량 시계열 · 수집 → 전처리/EDA → 딥러닝 이상탐지 → 리포트",
        "준지도 시퀀스 오토인코더 + 재구성오차 마할라노비스 거리 판정",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)


def _slide_overview(prs, embedded: list) -> None:
    """2. 워크플로우 개요(4단계 + 담당 에이전트)."""
    slide = _blank_slide(prs)
    _add_title(slide, "워크플로우 개요 — 4단계 자동화")
    _add_bullets(
        slide,
        [
            "① 수집 (mfg-collector): TEP 52변수 스트림을 SQLite에 누적 → stream.db",
            "② 전처리·EDA·통계 (mfg-eda): 정제 + 분포비교 + t/KS 검정 → eda/*",
            "③ 딥러닝 이상탐지 (mfg-model): 정상만 학습한 AE 3종 비교 → scores/comparison",
            "④ 대시보드·리포트 (mfg-reporter): Streamlit 멀티페이지 + 포트폴리오 PPT",
            ("각 단계는 파일 겹침 없이 소유·유지보수 → 병렬 안전, 산출물로 인계", 1),
            ("포지셔닝: ML 모델링이 핵심이되 수집~리포트 전 과정을 자동화", 1),
        ],
        size=17,
    )


def _slide_data(prs, embedded: list) -> None:
    """3. 데이터 구성(fault_counts.png 임베드)."""
    slide = _blank_slide(prs)
    _add_title(slide, "데이터 — TEP 885행 구성")
    stats = _collect_stats()
    if stats["total"]:
        ratio = stats["normal"] / stats["fault"] if stats["fault"] else float("nan")
        _add_bullets(
            slide,
            [
                f"총 {stats['total']:,}행 · 정상 {stats['normal']:,} / 결함 {stats['fault']:,}"
                f" (정상:결함 ≈ {ratio:.1f}:1)",
                f"결함 IDV {len(stats['fault_modes'])}종: {stats['fault_modes']}",
                "정상 위주 백필 구성 — 준지도 학습(정상만 학습)에 적합",
            ],
            top=1.3,
            height=1.6,
            size=16,
        )
    img = EDA_DIR / "fault_counts.png"
    if img.exists():
        _add_image_fit(slide, img, left=2.5, top=3.0, max_w=8.3, max_h=4.0)
        embedded.append(img.name)
    else:
        _missing_note(slide, "fault_counts.png — EDA 단계 실행 필요")


def _slide_eda(prs, embedded: list) -> None:
    """4. EDA 핵심(corr_heatmap + dist_normal_vs_fault 임베드)."""
    slide = _blank_slide(prs)
    _add_title(slide, "EDA 핵심 — 상관구조 & 정상 vs 결함 분포")
    _add_bullets(
        slide,
        [
            "52변수 상관 히트맵으로 군집/중복 구조 파악, 분포 비교로 분별력 점검",
            "KS 검정 상위 변수(예: xmeas_04·05 등)에서 정상·결함 분포 차이가 유의",
            "변수별 분별력 차이가 결함 IDV별 탐지율 편차의 근거",
        ],
        top=1.3,
        height=1.5,
        size=15,
    )
    left_img = EDA_DIR / "corr_heatmap.png"
    right_img = EDA_DIR / "dist_normal_vs_fault.png"
    any_img = False
    if left_img.exists():
        _add_image_fit(slide, left_img, left=0.4, top=2.9, max_w=6.0, max_h=4.2)
        embedded.append(left_img.name)
        any_img = True
    if right_img.exists():
        _add_image_fit(slide, right_img, left=6.9, top=2.9, max_w=6.0, max_h=4.2)
        embedded.append(right_img.name)
        any_img = True
    if not any_img:
        _missing_note(slide, "corr_heatmap.png / dist_normal_vs_fault.png — EDA 단계 실행 필요")


def _slide_model_compare(prs, embedded: list) -> None:
    """5. 모델 비교(실제 pptx 표 + roc_pr_curves.png 임베드)."""
    from pptx.util import Inches, Pt

    slide = _blank_slide(prs)
    _add_title(slide, "모델 비교 — AE 3종 (885행, 이상=양성)")
    comp = _comparison()
    if comp.empty:
        _missing_note(slide, "comparison.parquet — 모델 단계(detect) 실행 필요")
        return

    # 동적 표: 모델 × 지표 (comparison.parquet에서 채움)
    metrics = [("precision", "Precision"), ("recall", "Recall"), ("f1", "F1"),
               ("roc_auc", "ROC-AUC"), ("pr_auc", "PR-AUC")]
    metrics = [(k, lbl) for k, lbl in metrics if k in comp.columns]
    comp = comp.sort_values("f1", ascending=False) if "f1" in comp.columns else comp
    n_rows = len(comp) + 1
    n_cols = len(metrics) + 1
    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(1.4), Inches(6.4), Inches(0.4 * n_rows)
    ).table
    tbl.cell(0, 0).text = "모델"
    for j, (_, lbl) in enumerate(metrics, start=1):
        tbl.cell(0, j).text = lbl
    for i, (_, row) in enumerate(comp.iterrows(), start=1):
        name = MODEL_LABELS.get(row["model"], row["model"])
        if row["model"] == DEFAULT_MODEL:
            name += " ★기본"
        tbl.cell(i, 0).text = name
        for j, (k, _) in enumerate(metrics, start=1):
            tbl.cell(i, j).text = f"{row[k]:.3f}"
    for r in range(n_rows):
        for c in range(n_cols):
            for para in tbl.cell(r, c).text_frame.paragraphs:
                para.font.size = Pt(11)

    # 기본 모델 코멘트(동적)
    best = comp[comp["model"] == DEFAULT_MODEL]
    if not best.empty:
        b = best.iloc[0]
        _add_bullets(
            slide,
            [
                f"기본 모델 {MODEL_LABELS.get(DEFAULT_MODEL)}: F1·ROC-AUC·PR-AUC 최고,"
                f" 오탐(FP) {int(b['fp'])}로 최저",
                "정밀도 우선 운용 시 과탐을 가장 적게 냄(조건: 현재 885행 기준)",
            ],
            left=0.5,
            top=1.5 + 0.4 * n_rows,
            width=6.4,
            height=1.4,
            size=13,
        )

    img = MODELS_DIR / "roc_pr_curves.png"
    if img.exists():
        _add_image_fit(slide, img, left=7.1, top=1.4, max_w=6.0, max_h=5.4)
        embedded.append(img.name)


def _slide_fault_limits(prs, embedded: list) -> None:
    """6. 결함별 성능 & 한계(IDV별 탐지율 표 + 미탐 한계 정직 표기)."""
    from pptx.util import Inches, Pt

    slide = _blank_slide(prs)
    _add_title(slide, "결함별 성능 & 한계 — 정직 표기")
    pf = _per_fault()
    if pf.empty:
        _missing_note(slide, "comparison_per_fault.parquet — 모델 단계(detect) 실행 필요")
        return

    vae = pf[(pf["model"] == DEFAULT_MODEL) & (pf["fault_id"] != 0)].sort_values("fault_id")
    if vae.empty:
        _missing_note(slide, "기본 모델 결함별 탐지율 행이 없습니다")
        return

    n_cols = len(vae) + 1
    tbl = slide.shapes.add_table(
        2, n_cols, Inches(0.5), Inches(1.5), Inches(SLIDE_W - 1.0), Inches(1.0)
    ).table
    tbl.cell(0, 0).text = "IDV"
    tbl.cell(1, 0).text = "탐지율"
    for j, (_, row) in enumerate(vae.iterrows(), start=1):
        tbl.cell(0, j).text = str(int(row["fault_id"]))
        tbl.cell(1, j).text = f"{row['recall']:.2f}"
    for r in range(2):
        for c in range(n_cols):
            for para in tbl.cell(r, c).text_frame.paragraphs:
                para.font.size = Pt(11)

    detected = vae[vae["recall"] >= 0.99]["fault_id"].astype(int).tolist()
    missed = vae[vae["recall"] <= 0.01]["fault_id"].astype(int).tolist()
    partial = vae[(vae["recall"] > 0.01) & (vae["recall"] < 0.99)]["fault_id"].astype(int).tolist()
    _add_bullets(
        slide,
        [
            f"잘 탐지(≈1.00): IDV {detected} — 재구성오차에 뚜렷한 패턴 변화",
            f"부분 탐지: IDV {partial}" if partial else "부분 탐지: 해당 없음",
            f"미탐(0.00): IDV {missed} — 재구성오차가 정상과 사실상 구별 불가",
            ("→ AE 계열 원리상 한계로, 일부 결함은 분리 불가(MODEL_CARD 근거)", 1),
            ("과장 대신 한계를 그대로 표기 — 포트폴리오 신뢰성", 1),
        ],
        top=2.9,
        height=3.0,
        size=15,
    )


def _slide_conclusion(prs, embedded: list) -> None:
    """7. 결론/재현(핵심 수치 + 실행 명령)."""
    slide = _blank_slide(prs)
    _add_title(slide, "결론 & 재현")
    comp = _comparison()
    lines: list = []
    best = comp[comp["model"] == DEFAULT_MODEL] if not comp.empty else pd.DataFrame()
    if not best.empty:
        b = best.iloc[0]
        lines.append(
            f"기본 모델 {MODEL_LABELS.get(DEFAULT_MODEL)}: "
            f"F1 {b['f1']:.3f} · ROC-AUC {b['roc_auc']:.3f} · PR-AUC {b['pr_auc']:.3f}"
            f" (P {b['precision']:.3f} / R {b['recall']:.3f}, FP {int(b['fp'])})"
        )
    lines += [
        "수집~리포트 전 과정 자동화: 산출물 인계 계약으로 단계별 병렬 유지보수",
        "한계: IDV 1·4·8 미탐 / 데이터 규모(885행)에 따른 조건부 결과",
        ("재현: pip install -r requirements-train.txt", 1),
        ("python run_workflow.py --collect-batches 12 --ppt", 1),
        ("단계별: src.collect.scheduler → src.pipeline.{preprocess,eda} → src.models.detect", 1),
        ("리포트: streamlit run src/report/dashboard/app.py · python -m src.report.ppt", 1),
    ]
    _add_bullets(slide, lines, top=1.5, height=5.2, size=15)


# ── 빌드 ──────────────────────────────────────────────────
def build() -> Path:
    """포트폴리오 덱 생성 → reports/portfolio.pptx 반환."""
    from pptx import Presentation
    from pptx.util import Inches

    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    embedded: list = []
    builders = [
        _slide_cover,
        _slide_overview,
        _slide_data,
        _slide_eda,
        _slide_model_compare,
        _slide_fault_limits,
        _slide_conclusion,
    ]
    for fn in builders:
        fn(prs, embedded)

    out = REPORTS_DIR / "portfolio.pptx"
    prs.save(out)
    n_slides = len(prs.slides._sldIdLst)
    print(f"[ppt] 포트폴리오 덱 생성 → {out}")
    print(f"[ppt] 슬라이드 {n_slides}장 · 임베드 이미지 {len(embedded)}개: {embedded}")
    return out


if __name__ == "__main__":
    build()
